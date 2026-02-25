"""
文件内容提取服务
支持：PDF、Word (.docx/.doc)、PowerPoint (.pptx/.ppt)、纯文本 (.txt/.md)、Excel (.xlsx/.xls)
"""
import os
from typing import List, Tuple
from app.utils.logger import logger


def extract_text_from_pdf(file_path: str) -> str:
    """使用 pdfplumber 提取 PDF 文本"""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    texts.append(f"[第{page_num}页]\n{text.strip()}")
        return "\n\n".join(texts)
    except ImportError:
        raise ImportError("pdfplumber 未安装，请运行: pip install pdfplumber")
    except Exception as e:
        logger.error(f"PDF 解析失败 {file_path}: {e}")
        raise


def extract_text_from_docx(file_path: str) -> str:
    """使用 python-docx 提取 Word 文本"""
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))
        return "\n".join(texts)
    except ImportError:
        raise ImportError("python-docx 未安装，请运行: pip install python-docx")
    except Exception as e:
        logger.error(f"Word 解析失败 {file_path}: {e}")
        raise


def extract_text_from_pptx(file_path: str) -> str:
    """使用 python-pptx 提取 PowerPoint 文本"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"[第{slide_num}页]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except ImportError:
        raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")
    except Exception as e:
        logger.error(f"PPT 解析失败 {file_path}: {e}")
        raise


def extract_text_from_txt(file_path: str) -> str:
    """提取纯文本文件内容"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error(f"TXT 解析失败 {file_path}: {e}")
        raise


def extract_text_from_excel(file_path: str) -> str:
    """
    提取 Excel 文件内容（.xlsx / .xls）
    每个 Sheet 单独标注，每行格式化为 「列名=值」，方便 RAG 检索。
    空行、全空列自动跳过。
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".xlsx":
            return _extract_xlsx(file_path)
        elif ext == ".xls":
            return _extract_xls(file_path)
        else:
            raise ValueError(f"不支持的 Excel 格式: {ext}")
    except (ImportError, ModuleNotFoundError) as e:
        raise ImportError(str(e))
    except Exception as e:
        logger.error(f"Excel 解析失败 {file_path}: {e}")
        raise


def _extract_xlsx(file_path: str) -> str:
    """使用 openpyxl 解析 .xlsx"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl 未安装，请运行: pip install openpyxl")

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    all_sheets: List[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        sheet_text = _rows_to_text(sheet_name, rows)
        if sheet_text:
            all_sheets.append(sheet_text)

    wb.close()
    return "\n\n".join(all_sheets)


def _extract_xls(file_path: str) -> str:
    """使用 xlrd 解析 .xls（旧版 Excel）"""
    try:
        import xlrd
    except ImportError:
        raise ImportError("xlrd 未安装，请运行: pip install xlrd")

    wb = xlrd.open_workbook(file_path)
    all_sheets: List[str] = []

    for sheet in wb.sheets():
        rows = [
            tuple(sheet.cell_value(r, c) for c in range(sheet.ncols))
            for r in range(sheet.nrows)
        ]
        sheet_text = _rows_to_text(sheet.name, rows)
        if sheet_text:
            all_sheets.append(sheet_text)

    return "\n\n".join(all_sheets)


def _rows_to_text(sheet_name: str, rows: List[tuple]) -> str:
    """
    将行列数据转换为可供 RAG 检索的文本。
    策略：
    - 首行作为表头（列名）
    - 每个数据行格式化为：「列名1=值1 | 列名2=值2 | …」
    - 若无表头则直接以 tab 分隔输出
    - 过滤空行
    """
    if not rows:
        return ""

    # 过滤全为空的行
    non_empty_rows = [
        r for r in rows
        if any(v is not None and str(v).strip() != "" for v in r)
    ]
    if not non_empty_rows:
        return ""

    lines: List[str] = [f"[Sheet: {sheet_name}]"]

    # 取第一行作为表头
    header = [str(v).strip() if v is not None else "" for v in non_empty_rows[0]]
    has_header = any(h for h in header)

    if has_header:
        # 输出表头行（作为列名参考）
        lines.append("列名：" + " | ".join(h for h in header if h))
        # 输出数据行
        for row in non_empty_rows[1:]:
            cells = []
            for col_idx, val in enumerate(row):
                col_name = header[col_idx] if col_idx < len(header) else f"列{col_idx + 1}"
                cell_str = str(val).strip() if val is not None else ""
                if cell_str and cell_str != "None":
                    cells.append(f"{col_name}={cell_str}")
            if cells:
                lines.append(" | ".join(cells))
    else:
        # 无表头，直接输出每行
        for row in non_empty_rows:
            cell_strs = [
                str(v).strip()
                for v in row
                if v is not None and str(v).strip() not in ("", "None")
            ]
            if cell_strs:
                lines.append("\t".join(cell_strs))

    # 只有标题行，没有实际内容
    if len(lines) <= 2:
        return ""

    return "\n".join(lines)


# 文件扩展名 → 解析函数映射
EXTRACTORS = {
    ".pdf":  extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc":  extract_text_from_docx,   # 旧版 Word 也尝试 python-docx
    ".pptx": extract_text_from_pptx,
    ".ppt":  extract_text_from_pptx,   # 旧版 PPT
    ".txt":  extract_text_from_txt,
    ".md":   extract_text_from_txt,
    ".xlsx": extract_text_from_excel,  # Excel 新版
    ".xls":  extract_text_from_excel,  # Excel 旧版
}

SUPPORTED_EXTENSIONS = set(EXTRACTORS.keys())


def extract_text_from_file(file_path: str) -> str:
    """
    根据文件扩展名自动选择解析器提取文本
    """
    ext = os.path.splitext(file_path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(
            f"不支持的文件格式: {ext}。"
            f"支持的格式：{', '.join(SUPPORTED_EXTENSIONS)}"
        )
    logger.info(f"正在提取文件内容: {file_path} (格式: {ext})")
    text = extractor(file_path)
    logger.info(f"文件提取完成: {file_path}，字符数={len(text)}")
    return text


def extract_texts_from_dir(dir_path: str) -> List[dict]:
    """
    遍历目录，提取所有支持格式文件的文本
    返回：[{"filename": str, "filepath": str, "text": str}]
    """
    results = []
    if not os.path.isdir(dir_path):
        return results

    for filename in os.listdir(dir_path):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        filepath = os.path.join(dir_path, filename)
        try:
            text = extract_text_from_file(filepath)
            if text.strip():
                results.append({
                    "filename": filename,
                    "filepath": filepath,
                    "text": text,
                })
        except Exception as e:
            logger.error(f"跳过文件 {filename}: {e}")
            continue

    return results


# ─── 网络 URL 提取 ─────────────────────────────────────────────────────────────

def _is_file_url(url: str, content_type: str = None) -> Tuple[bool, str]:
    """
    检测 URL 是否指向文件（而非 HTML 网页）
    返回：(is_file, file_extension)
    """
    # 从 URL 中提取文件扩展名（去除查询参数）
    url_path = url.split("?")[0].split("#")[0]
    url_ext = os.path.splitext(url_path)[1].lower()
    
    # 支持的文件扩展名
    file_extensions = {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".txt", ".md"}
    
    # 方法1：从 URL 扩展名判断
    if url_ext in file_extensions:
        return True, url_ext
    
    # 方法2：从 Content-Type 判断
    if content_type:
        content_type_lower = content_type.lower()
        # 常见的文件 MIME 类型
        file_mime_types = {
            "application/pdf": ".pdf",
            "application/msword": ".doc",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
            "application/vnd.ms-excel": ".xls",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
            "application/vnd.ms-powerpoint": ".ppt",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
            "text/plain": ".txt",
            "text/markdown": ".md",
        }
        for mime_type, ext in file_mime_types.items():
            if mime_type in content_type_lower:
                return True, ext
    
    return False, ""


def extract_text_from_url(url: str, timeout: int = 30) -> dict:
    """
    从网络 URL 抓取并提取文本内容。
    支持两种类型：
    1. 文件类型（PDF、Word、Excel、PPT等）：下载文件并提取文本
    2. HTML 网页：解析 HTML 提取文本
    
    返回：{"url": str, "title": str, "text": str}

    注意：对于需要 JavaScript 渲染（如腾讯文档、Google Docs 等）的页面，
    此方法仅能获取页面的静态 HTML 内容，可能无法获取完整数据。
    """
    try:
        import requests
        from bs4 import BeautifulSoup
        import tempfile
    except ImportError:
        raise ImportError("requests 或 beautifulsoup4 未安装，请运行: pip install requests beautifulsoup4 lxml")

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        ),
        "Accept": "*/*",  # 接受所有类型，包括文件
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    }

    logger.info(f"正在抓取 URL: {url}")
    
    # 先基于 URL 扩展名判断是否为文件
    is_file, file_ext = _is_file_url(url, None)
    
    # 发送 GET 请求获取响应（使用 stream=True 以便处理大文件）
    resp = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True, stream=True)
    resp.raise_for_status()
    
    # 从响应头获取 Content-Type，再次确认是否为文件
    content_type = resp.headers.get("content-type", "")
    if not is_file:
        is_file, file_ext = _is_file_url(url, content_type)
    
    if is_file:
        # 处理文件类型：下载并提取文本
        logger.info(f"检测到文件类型 URL: {url} (扩展名: {file_ext})")
        
        # 从 Content-Disposition 或 URL 中提取文件名
        filename = None
        content_disposition = resp.headers.get("content-disposition", "")
        if content_disposition:
            import re
            match = re.search(r'filename[^;=\n]*=(([\'"]).*?\2|[^\s;]+)', content_disposition)
            if match:
                filename = match.group(1).strip('"\'')
        
        if not filename:
            # 从 URL 提取文件名
            url_path = url.split("?")[0].split("#")[0]
            filename = os.path.basename(url_path) or f"file{file_ext}"
        
        # 确保文件名有正确的扩展名
        if not filename.endswith(file_ext):
            filename = filename.rsplit(".", 1)[0] + file_ext
        
        # 保存到临时文件
        temp_dir = tempfile.gettempdir()
        temp_file_path = os.path.join(temp_dir, f"zerag_download_{os.urandom(8).hex()}{file_ext}")
        
        try:
            with open(temp_file_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
            
            logger.info(f"文件已下载到临时路径: {temp_file_path}")
            
            # 使用现有的文件提取函数提取文本
            full_text = extract_text_from_file(temp_file_path)
            title = filename
            
            logger.info(f"文件提取完成: {url}，文件名={filename!r}，字符数={len(full_text)}")
            return {"url": url, "title": title, "text": full_text}
            
        finally:
            # 清理临时文件
            try:
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
                    logger.debug(f"已删除临时文件: {temp_file_path}")
            except Exception as e:
                logger.warning(f"删除临时文件失败 {temp_file_path}: {e}")
    
    # 处理 HTML 网页类型
    # 由于使用了 stream=True，响应内容还未被读取，可以直接使用 resp.text
    # 但需要先关闭流模式或使用 resp.content
    # 对于 HTML，我们重新发送请求（不使用 stream）以获取完整内容
    resp_html = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True)
    resp_html.raise_for_status()

    # 自动检测编码
    content_type_html = resp_html.headers.get("content-type", "")
    if "charset=" in content_type_html:
        encoding = content_type_html.split("charset=")[-1].split(";")[0].strip()
        resp_html.encoding = encoding
    else:
        resp_html.encoding = resp_html.apparent_encoding or "utf-8"

    soup = BeautifulSoup(resp_html.text, "lxml")

    # 提取标题
    title_tag = soup.find("title")
    title = title_tag.get_text(strip=True) if title_tag else url

    # 移除干扰元素
    for tag in soup(["script", "style", "noscript", "nav", "footer",
                     "header", "aside", "iframe", "svg", "button", "form"]):
        tag.decompose()

    # 优先提取正文区域
    main_content = (
        soup.find("main") or
        soup.find("article") or
        soup.find(id="content") or
        soup.find(class_="content") or
        soup.find(id="main") or
        soup.body
    )
    if main_content is None:
        main_content = soup

    # 提取文本，保留段落结构
    lines = []
    for element in main_content.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6",
                                           "li", "td", "th", "div", "span"]):
        text = element.get_text(separator=" ", strip=True)
        if text and len(text) > 5:  # 过滤极短噪音
            lines.append(text)

    # 去重并合并
    seen = set()
    unique_lines = []
    for line in lines:
        if line not in seen:
            seen.add(line)
            unique_lines.append(line)

    full_text = "\n".join(unique_lines)

    # 如果正文提取为空，回退到全文本
    if not full_text.strip():
        full_text = soup.get_text(separator="\n", strip=True)

    logger.info(f"URL 抓取完成: {url}，标题={title!r}，字符数={len(full_text)}")
    return {"url": url, "title": title, "text": full_text}


def extract_texts_from_urls(urls: list) -> list:
    """
    批量抓取 URL 列表，返回：[{"url": str, "title": str, "text": str}]
    单个 URL 失败不影响其他 URL。
    """
    results = []
    for url in urls:
        try:
            doc = extract_text_from_url(url)
            if doc["text"].strip():
                results.append(doc)
            else:
                logger.warning(f"URL {url} 提取内容为空，已跳过")
        except Exception as e:
            logger.error(f"跳过 URL {url}: {e}")
    return results
